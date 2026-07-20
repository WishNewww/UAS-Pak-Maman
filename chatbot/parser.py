"""
parser.py
=========
Intelligent PowerPoint document parser for the Regulasi Meteorologi Chatbot.

Extracts structured content from .pptx files including:
- Slide titles and subtitles
- Bullet points with hierarchy preservation
- Paragraphs and body text
- Speaker notes
- Tables (all cells, row-by-row)
- Numbered lists
- Section headers
- Captions

Decorative / layout placeholders (background images, logos, theme shapes)
are automatically filtered out.

Metadata per slide:
    folder, filename, slide_number, slide_title, topic, section

Author : Regulasi Meteorologi Chatbot
Version: 1.0.0
"""

from __future__ import annotations

import logging
import re
from dataclasses import dataclass, field
from pathlib import Path
from typing import Generator

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt

from config import AppConfig, DataConfig

logger = logging.getLogger("meteorologi_chatbot.parser")


# ---------------------------------------------------------------------------
# Data Structures
# ---------------------------------------------------------------------------

@dataclass
class SlideContent:
    """Structured content extracted from a single slide."""

    folder: str
    filename: str
    file_path: str
    slide_number: int
    slide_title: str
    subtitle: str
    section: str
    topic: str
    body_lines: list[str] = field(default_factory=list)
    notes: str = ""
    tables: list[list[list[str]]] = field(default_factory=list)

    def to_text_block(self) -> str:
        """
        Render the slide as a single coherent text block.
        Preserves hierarchy using indentation markers.
        """
        parts: list[str] = []

        if self.slide_title:
            parts.append(f"[Judul Slide]: {self.slide_title}")
        if self.subtitle:
            parts.append(f"[Subjudul]: {self.subtitle}")
        if self.section:
            parts.append(f"[Bagian/Section]: {self.section}")

        if self.body_lines:
            parts.append("[Konten]:")
            parts.extend(self.body_lines)

        for idx, table in enumerate(self.tables, start=1):
            parts.append(f"[Tabel {idx}]:")
            for row in table:
                parts.append(" | ".join(cell for cell in row if cell))

        if self.notes:
            parts.append(f"[Catatan Pembicara]: {self.notes}")

        return "\n".join(parts)

    @property
    def metadata(self) -> dict:
        return {
            "folder": self.folder,
            "filename": self.filename,
            "file_path": self.file_path,
            "slide_number": self.slide_number,
            "slide_title": self.slide_title,
            "subtitle": self.subtitle,
            "section": self.section,
            "topic": self.topic,
        }


# ---------------------------------------------------------------------------
# Internal helpers
# ---------------------------------------------------------------------------

# Placeholder types that are pure decoration / not content
_DECORATIVE_PLACEHOLDER_TYPES = {
    "PICTURE",
    "MEDIA",
    "OBJECT",
    "THUMBNAIL",
    "SLIDE_IMAGE",
}

# Regex for cleaning excess whitespace while preserving newlines
_MULTI_SPACE = re.compile(r"[ \t]{2,}")
_TRAILING_WS = re.compile(r"[ \t]+$", re.MULTILINE)


def _clean_text(text: str) -> str:
    """Remove redundant spaces / tabs; preserve line breaks and Indonesian text."""
    if not text:
        return ""
    text = _MULTI_SPACE.sub(" ", text)
    text = _TRAILING_WS.sub("", text)
    text = text.strip()
    return text


def _is_decorative(shape) -> bool:
    """Return True for shapes that carry no meaningful textual content."""
    if not shape.has_text_frame and not hasattr(shape, "table"):
        return True
    # Filter out picture/media placeholders
    if shape.is_placeholder:
        ph_type = str(shape.placeholder_format.type).split(".")[-1].upper()
        if ph_type in _DECORATIVE_PLACEHOLDER_TYPES:
            return True
    return False


def _extract_text_frame_lines(text_frame, base_indent: int = 0) -> list[str]:
    """
    Extract lines from a text frame preserving bullet hierarchy.

    - Level 0 → no indent prefix
    - Level 1 → "  • "
    - Level 2 → "    ◦ "
    - Level 3+ → "      ‣ "
    """
    lines: list[str] = []
    indent_map = {0: "", 1: "  • ", 2: "    ◦ ", 3: "      ‣ "}

    for para in text_frame.paragraphs:
        raw = para.text.strip()
        if not raw:
            continue

        raw = _clean_text(raw)
        level = min((para.level or 0) + base_indent, 3)
        prefix = indent_map.get(level, "      ‣ ")
        lines.append(f"{prefix}{raw}")

    return lines


def _extract_notes(slide) -> str:
    """Return cleaned speaker notes text, or empty string."""
    try:
        notes_slide = slide.notes_slide
        tf = notes_slide.notes_text_frame
        lines = []
        for para in tf.paragraphs:
            text = _clean_text(para.text)
            if text:
                lines.append(text)
        return " ".join(lines)
    except Exception:
        return ""


def _extract_table(shape) -> list[list[str]]:
    """Extract table contents as a 2-D list of strings."""
    rows_data: list[list[str]] = []
    try:
        table = shape.table
        for row in table.rows:
            row_data = [_clean_text(cell.text) for cell in row.cells]
            # Only include rows that have at least one non-empty cell
            if any(cell for cell in row_data):
                rows_data.append(row_data)
    except Exception as exc:
        logger.debug("Table extraction error: %s", exc)
    return rows_data


def _derive_topic(filename: str) -> str:
    """
    Infer a topic label from the filename.
    E.g. 'Regulasi Internasional Met 09_M8C_250526.pptx' → 'Materi 09'
    """
    stem = Path(filename).stem
    # Match patterns like 'Met 09', 'Met09', 'Materi 09'
    match = re.search(r"(?:Met\s*|Materi\s*)(\d+)", stem, re.IGNORECASE)
    if match:
        return f"Materi {int(match.group(1)):02d}"
    return stem


# ---------------------------------------------------------------------------
# Core Parser
# ---------------------------------------------------------------------------

class PowerPointParser:
    """
    Recursively scans a directory for .pptx files and extracts
    structured slide content.
    """

    def __init__(self, data_dir: Path | None = None) -> None:
        self.data_dir: Path = data_dir or DataConfig.DATA_DIR
        logger.info("PowerPointParser initialised. Data dir: %s", self.data_dir)

    # ------------------------------------------------------------------
    # Public API
    # ------------------------------------------------------------------

    def discover_files(self) -> list[Path]:
        """Return all .pptx files found recursively under data_dir."""
        if not self.data_dir.exists():
            raise FileNotFoundError(f"Data directory not found: {self.data_dir}")

        files = sorted(self.data_dir.rglob("*.pptx"))
        logger.info("Discovered %d .pptx file(s) in '%s'", len(files), self.data_dir)
        return files

    def parse_file(self, file_path: Path) -> list[SlideContent]:
        """
        Parse a single .pptx file and return a list of SlideContent objects,
        one per slide.
        """
        logger.debug("Parsing file: %s", file_path)
        slides: list[SlideContent] = []

        try:
            prs = Presentation(str(file_path))
        except Exception as exc:
            logger.error("Failed to open '%s': %s", file_path, exc)
            return slides

        folder_name = file_path.parent.name
        filename = file_path.name
        topic = _derive_topic(filename)
        current_section: str = ""

        for slide_idx, slide in enumerate(prs.slides, start=1):
            content = self._parse_slide(
                slide=slide,
                slide_number=slide_idx,
                folder=folder_name,
                filename=filename,
                file_path=str(file_path),
                topic=topic,
                current_section=current_section,
            )
            # Update running section tracker
            if content.section:
                current_section = content.section

            slides.append(content)

        logger.info(
            "Parsed '%s': %d slide(s) extracted.", filename, len(slides)
        )
        return slides

    def parse_all(self) -> list[SlideContent]:
        """Parse every .pptx file discovered under data_dir."""
        all_slides: list[SlideContent] = []
        files = self.discover_files()

        for fp in files:
            slide_contents = self.parse_file(fp)
            all_slides.extend(slide_contents)

        logger.info(
            "Total slides extracted from all files: %d", len(all_slides)
        )
        return all_slides

    def iter_parse(self) -> Generator[SlideContent, None, None]:
        """Lazy generator version of parse_all — memory-efficient for large sets."""
        for fp in self.discover_files():
            yield from self.parse_file(fp)

    # ------------------------------------------------------------------
    # Private: single-slide extraction
    # ------------------------------------------------------------------

    def _parse_slide(
        self,
        slide,
        slide_number: int,
        folder: str,
        filename: str,
        file_path: str,
        topic: str,
        current_section: str,
    ) -> SlideContent:
        slide_title = ""
        subtitle = ""
        body_lines: list[str] = []
        tables: list[list[list[str]]] = []
        section = current_section

        for shape in slide.shapes:
            if _is_decorative(shape):
                continue

            # ---- Tables ------------------------------------------------
            if shape.has_table:
                table_data = _extract_table(shape)
                if table_data:
                    tables.append(table_data)
                continue

            # ---- Text frames -------------------------------------------
            if not shape.has_text_frame:
                continue

            # Identify placeholder role
            if shape.is_placeholder:
                ph_idx = shape.placeholder_format.idx
                ph_type = str(shape.placeholder_format.type).split(".")[-1].upper()

                if ph_idx == 0 or ph_type in ("TITLE", "CENTER_TITLE"):
                    # Primary title
                    slide_title = _clean_text(shape.text_frame.text)
                    # A section heading is typically a title-only slide or
                    # a shape explicitly marked as a section header
                    if self._looks_like_section(slide_title):
                        section = slide_title

                elif ph_idx == 1 and ph_type in (
                    "SUBTITLE", "BODY", "CENTER_BODY"
                ):
                    raw_lines = _extract_text_frame_lines(shape.text_frame)
                    if ph_type == "SUBTITLE":
                        subtitle = " ".join(raw_lines)
                    else:
                        body_lines.extend(raw_lines)

                else:
                    # Other content placeholders
                    body_lines.extend(
                        _extract_text_frame_lines(shape.text_frame)
                    )
            else:
                # Non-placeholder text boxes
                lines = _extract_text_frame_lines(shape.text_frame)
                if lines:
                    body_lines.extend(lines)

        # Deduplicate body lines while preserving order
        seen: set[str] = set()
        deduped: list[str] = []
        for line in body_lines:
            if line not in seen:
                seen.add(line)
                deduped.append(line)

        notes = _extract_notes(slide)

        return SlideContent(
            folder=folder,
            filename=filename,
            file_path=file_path,
            slide_number=slide_number,
            slide_title=slide_title,
            subtitle=subtitle,
            section=section,
            topic=topic,
            body_lines=deduped,
            notes=notes,
            tables=tables,
        )

    @staticmethod
    def _looks_like_section(title: str) -> bool:
        """
        Heuristic: a slide title is a section marker if it is short
        and does not look like a regular sentence.
        """
        if not title:
            return False
        words = title.split()
        # Short titles (≤ 6 words) with no period → likely a section heading
        return len(words) <= 6 and not title.endswith(".")


# ---------------------------------------------------------------------------
# Quick-test entry point
# ---------------------------------------------------------------------------

if __name__ == "__main__":
    import json

    logging.basicConfig(level=logging.DEBUG)
    parser = PowerPointParser()

    try:
        all_slides = parser.parse_all()
        print(f"\n✅ Total slides parsed: {len(all_slides)}\n")

        # Print a preview of the first 3 slides
        for slide in all_slides[:3]:
            print("=" * 60)
            print(f"File   : {slide.filename}")
            print(f"Slide  : {slide.slide_number}")
            print(f"Title  : {slide.slide_title}")
            print(f"Topic  : {slide.topic}")
            print(f"Section: {slide.section}")
            print("Content preview:")
            print(slide.to_text_block()[:500])
            print()
    except FileNotFoundError as e:
        print(f"❌ {e}")
